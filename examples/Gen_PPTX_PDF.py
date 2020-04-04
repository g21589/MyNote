from pptx import Presentation
from pptx.util import Inches, Cm, Pt

from fpdf import FPDF

########
# PPTX #
########
prs = Presentation('template.pptx')

prs.slide_width = Cm(33.876)
prs.slide_height = Cm(19.05)

slide = prs.slides.add_slide(prs.slide_layouts[1])

slide.shapes.title.text = 'Test'

pic = slide.shapes.add_picture('plot.png', 0, prs.slide_height/12*4, width=prs.slide_width, height=prs.slide_height/12*3)
slide.shapes[0]._element.addprevious(pic._element)

pic = slide.shapes.add_picture('plot.png', 0, prs.slide_height/12*8, width=prs.slide_width, height=prs.slide_height/12*3)
slide.shapes[1]._element.addprevious(pic._element)

prs.save('Result.pptx')

#######
# PDF #
#######
pdf = FPDF()
pdf.add_page()
pdf.image('plot.png', x=1, y=1, w=200)
pdf.set_font("Arial", size=12)
pdf.ln(85)  # move 85 down
pdf.cell(200, 10, txt="{}".format('plot.png'), ln=1)
pdf.output("Result.pdf")
